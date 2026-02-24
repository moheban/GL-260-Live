import os, json, sys
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.ticker import (
    MultipleLocator,
    AutoMinorLocator,
    AutoLocator,
    FuncFormatter,
    ScalarFormatter,
)
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg, NavigationToolbar2Tk
from tkinter import colorchooser, simpledialog

# Optional (hover tooltips)
try:
    import mplcursors  # pip install mplcursors
except Exception:
    mplcursors = None

# Optional (PowerPoint export)
try:
    from pptx import Presentation  # pip install python-pptx
    from pptx.util import Inches
except Exception:
    Presentation = None

# ---------- Global styling ----------
plt.rcParams["font.family"] = "serif"
plt.rcParams["font.serif"] = ["Times New Roman"] + plt.rcParams["font.serif"]

SUBPLOT_TITLE_FONTSIZE = 14
LABEL_FONTSIZE = 12
SUPTITLE_FONTSIZE = 18
SUPTITLE_Y = 0.975
TICK_LABELSIZE = 11

LINESTYLE_DEFAULT = "-"
MARKER_DEFAULT = "."
LINEWIDTH_DEFAULT = 1.0
MARKERSIZE_DEFAULT = 2.0

SETTINGS_FILE = "general_plotter_settings.json"
MARKER_CHOICES = [
    "",
    ".",
    "o",
    "v",
    "^",
    "<",
    ">",
    "s",
    "p",
    "*",
    "x",
    "D",
    "h",
    "+",
    "1",
    "2",
    "3",
    "4",
]

# Resizable-left-column defaults
LEFT_COL_INIT = 520  # initial width of the left pane (drag the sash to resize)
LEFT_COL_MIN = 360  # don't let it shrink below this (keeps controls readable)
SCROLLBAR_FALLBACK_WIDTH = 18  # used only before the real width is known

EXPORT_DPI = 600


def _safe_float(v, default=None):
    try:
        return float(v)
    except Exception:
        return default


def _format_thousands(v, _):
    try:
        v = float(v)
    except Exception:
        return ""
    abs_v = abs(v)
    if abs_v >= 1_000_000_000:
        return f"{v/1_000_000_000:.2f}B"
    if abs_v >= 1_000_000:
        return f"{v/1_000_000:.2f}M"
    if abs_v >= 1_000:
        return f"{v/1_000:.2f}k"
    return f"{v:.0f}"


def _percent_fmt(v, _):
    try:
        return f"{100*float(v):.0f}%"
    except Exception:
        return ""


def _make_legend_draggable(legend):
    if legend is None:
        return
    try:
        legend.set_draggable(True)
    except Exception:
        pass


def _make_legends_draggable(fig):
    if fig is None:
        return
    seen = set()
    for legend in getattr(fig, "legends", []):
        if legend is None:
            continue
        _make_legend_draggable(legend)
        seen.add(id(legend))
    for ax in getattr(fig, "axes", []):
        legend = ax.get_legend()
        if legend is None or id(legend) in seen:
            continue
        _make_legend_draggable(legend)


def _resolve_right_label(custom, fallback):
    custom = (custom or "").strip()
    if custom:
        return custom
    return fallback or ""


def _apply_theme(theme_key: str):
    if theme_key == "classic":
        plt.rcParams.update(
            {
                "axes.grid": False,
                "figure.facecolor": "white",
                "axes.facecolor": "white",
                "text.color": "black",
                "axes.labelcolor": "black",
                "xtick.color": "black",
                "ytick.color": "black",
            }
        )
    elif theme_key == "minimal":
        plt.rcParams.update(
            {
                "axes.grid": True,
                "grid.alpha": 0.25,
                "figure.facecolor": "white",
                "axes.facecolor": "white",
                "text.color": "black",
                "axes.labelcolor": "black",
                "xtick.color": "black",
                "ytick.color": "black",
            }
        )
    elif theme_key == "dark":
        plt.rcParams.update(
            {
                "axes.grid": True,
                "grid.color": "#888",
                "figure.facecolor": "#111",
                "axes.facecolor": "#111",
                "text.color": "white",
                "axes.labelcolor": "white",
                "xtick.color": "white",
                "ytick.color": "white",
            }
        )
    elif theme_key == "journal":
        plt.rcParams.update(
            {
                "axes.grid": True,
                "grid.linestyle": ":",
                "grid.alpha": 0.35,
                "figure.facecolor": "white",
                "axes.facecolor": "white",
                "text.color": "black",
                "axes.labelcolor": "black",
                "xtick.color": "black",
                "ytick.color": "black",
            }
        )


def _apply_formatter_to_axis(axis, fmt_key: str):
    if fmt_key == "plain":
        axis.set_major_formatter(ScalarFormatter(useMathText=False))
    elif fmt_key == "sci":
        sf = ScalarFormatter(useMathText=True)
        sf.set_powerlimits((-3, 3))
        axis.set_major_formatter(sf)
    elif fmt_key == "percent":
        axis.set_major_formatter(FuncFormatter(_percent_fmt))
    elif fmt_key == "thousands":
        axis.set_major_formatter(FuncFormatter(_format_thousands))


def _safe_color_dialog(initial="#1f77b4"):
    try:
        c = colorchooser.askcolor(color=initial)[1]
        return c or initial
    except Exception:
        return initial


class GeneralPlotter(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("General Plotter")
        self.minsize(1200, 840)
        self.geometry(self._load_window_geom() or "1260x920+50+50")

        # State
        self.df: pd.DataFrame | None = None
        self.file_path = ""
        self.sheet_names: list[str] = []
        self.current_sheet = tk.StringVar()
        self.columns: list[str] = []
        self.nb: ttk.Notebook | None = None
        self.render_tab_counter = 0
        self.render_tabs: list[dict] = []
        self._sheet_indicator_canvas: tk.Canvas | None = None
        self._sheet_indicator_state = False

        # Persistent settings
        self.settings = self._load_settings()

        # ---- Tk Vars (seed from settings) ----
        self.title_text = tk.StringVar(value=self.settings.get("title_text", ""))
        self.suptitle_text = tk.StringVar(value=self.settings.get("suptitle_text", ""))

        default_plot_type = self.settings.get("plot_type") or self._map_old_mode(
            self.settings.get("plot_mode")
        )
        self.plot_type = tk.StringVar(value=default_plot_type or "single")

        self.legend_loc = tk.StringVar(
            value=self.settings.get("legend_loc", "lower center")
        )
        self.plot_kind = tk.StringVar(value=self.settings.get("plot_kind", "line"))

        self.auto_x = tk.BooleanVar(value=self.settings.get("auto_x", True))
        self.auto_y_left = tk.BooleanVar(value=self.settings.get("auto_y_left", True))
        self.auto_y_right = tk.BooleanVar(value=self.settings.get("auto_y_right", True))

        self.x_min = tk.StringVar(value=str(self.settings.get("x_min", "")))
        self.x_max = tk.StringVar(value=str(self.settings.get("x_max", "")))
        self.y_left_min = tk.StringVar(value=str(self.settings.get("y_left_min", "")))
        self.y_left_max = tk.StringVar(value=str(self.settings.get("y_left_max", "")))
        self.y_right_min = tk.StringVar(value=str(self.settings.get("y_right_min", "")))
        self.y_right_max = tk.StringVar(value=str(self.settings.get("y_right_max", "")))

        self.auto_time_ticks = tk.BooleanVar(
            value=self.settings.get("auto_time_ticks", True)
        )
        self.auto_y_ticks = tk.BooleanVar(value=self.settings.get("auto_y_ticks", True))
        self.xmaj = tk.StringVar(value=str(self.settings.get("x_major_tick", "")))
        self.xminr = tk.StringVar(value=str(self.settings.get("x_minor_tick", "")))
        self.ymaj = tk.StringVar(value=str(self.settings.get("y_major_tick", "")))
        self.yminr = tk.StringVar(value=str(self.settings.get("y_minor_tick", "")))

        self.linewidth = tk.DoubleVar(
            value=self.settings.get("linewidth", LINEWIDTH_DEFAULT)
        )
        self.markersize = tk.DoubleVar(
            value=self.settings.get("markersize", MARKERSIZE_DEFAULT)
        )

        self.ncols = tk.IntVar(value=int(self.settings.get("ncols", 2)))

        self.right_axis_series = tk.StringVar(
            value=self.settings.get("right_axis_series", "None")
        )
        self.facet_col_saved = self.settings.get("facet_col", "None")

        self._preview_after_id = None
        self._preview_debounce_ms = 250

        self.series_style: dict[str, dict] = self.settings.get("series_style", {})
        self.references: list[dict] = self.settings.get("references", [])
        self.filter_query = tk.StringVar(value=self.settings.get("filter_query", ""))
        self.decimate_preview = tk.BooleanVar(
            value=self.settings.get("decimate_preview", True)
        )
        # NEW: custom axis label text
        self.x_label = tk.StringVar(value=self.settings.get("x_label", ""))
        self.y_left_label = tk.StringVar(value=self.settings.get("y_left_label", ""))
        self.y_right_label = tk.StringVar(value=self.settings.get("y_right_label", ""))
        self.yerr_col = tk.StringVar(value=self.settings.get("yerr_col", "None"))
        self.y_left_format = tk.StringVar(
            value=self.settings.get("y_left_format", "plain")
        )
        self.y_right_format = tk.StringVar(
            value=self.settings.get("y_right_format", "plain")
        )
        self.theme = tk.StringVar(value=self.settings.get("theme", "classic"))
        self.top_k_facets = tk.IntVar(value=int(self.settings.get("top_k_facets", 0)))
        self.show_grid = tk.BooleanVar(value=True)  # NEW: toggle for grid on/off

        # UI
        self._build_ui()

        # Reload last-used file if present
        last_path = self.settings.get("last_file_path", "")
        if last_path and os.path.exists(last_path):
            self._load_file(last_path, auto=True)

        self.protocol("WM_DELETE_WINDOW", self._on_close)

    def _collect_legend_entries(self, fig):
        """Collect unique legend handles/labels from ALL axes in a figure (including twinx)."""
        handles, labels, seen = [], [], set()
        for ax in fig.axes:
            h, l = ax.get_legend_handles_labels()
            for hi, li in zip(h, l):
                if not li or li in seen:
                    continue
                handles.append(hi)
                labels.append(li)
                seen.add(li)
        return handles, labels

    # ---------- Back-compat ----------
    @staticmethod
    def _map_old_mode(old_mode: str | None) -> str | None:
        if not old_mode:
            return None
        return {"single": "single", "by_y": "grid", "facet": "facet"}.get(
            old_mode, "single"
        )

    # ---------- Persistence ----------
    def _load_settings(self):
        try:
            if os.path.exists(SETTINGS_FILE):
                with open(SETTINGS_FILE, "r") as f:
                    return json.load(f)
        except Exception:
            pass
        return {}

    def _save_settings(self):
        payload = dict(self.settings)
        payload.update(
            {
                "title_text": self.title_text.get(),
                "suptitle_text": self.suptitle_text.get(),
                "plot_type": self.plot_type.get(),
                "legend_loc": self.legend_loc.get(),
                "plot_kind": self.plot_kind.get(),
                "auto_x": self.auto_x.get(),
                "auto_y_left": self.auto_y_left.get(),
                "auto_y_right": self.auto_y_right.get(),
                "x_min": self.x_min.get(),
                "x_max": self.x_max.get(),
                "y_left_min": self.y_left_min.get(),
                "y_left_max": self.y_left_max.get(),
                "y_right_min": self.y_right_min.get(),
                "y_right_max": self.y_right_max.get(),
                "auto_time_ticks": self.auto_time_ticks.get(),
                "auto_y_ticks": self.auto_y_ticks.get(),
                "x_major_tick": self.xmaj.get(),
                "x_minor_tick": self.xminr.get(),
                "y_major_tick": self.ymaj.get(),
                "y_minor_tick": self.yminr.get(),
                "linewidth": self.linewidth.get(),
                "markersize": self.markersize.get(),
                "ncols": self.ncols.get(),
                "x_col": self.cb_x.get() if hasattr(self, "cb_x") else "",
                "facet_col": (
                    self.cb_facet.get() if hasattr(self, "cb_facet") else "None"
                ),
                "y_cols": self._get_selected_y(),
                "right_axis_series": self.right_axis_series.get(),
                "log_x": self.log_x.get(),
                "log_y_left": self.log_y_left.get(),
                "log_y_right": self.log_y_right.get(),
                "x_tick_rotation": self.x_tick_rotation.get(),
                "series_style": self.series_style,
                "references": self.references,
                "filter_query": self.filter_query.get(),
                "decimate_preview": self.decimate_preview.get(),
                "x_label": self.x_label.get(),  # NEW
                "y_left_label": self.y_left_label.get(),  # NEW
                "y_right_label": self.y_right_label.get(),  # NEW
                "show_grid": self.show_grid.get(),
                "yerr_col": self.yerr_col.get(),
                "y_left_format": self.y_left_format.get(),
                "y_right_format": self.y_right_format.get(),
                "theme": self.theme.get(),
                "top_k_facets": self.top_k_facets.get(),
            }
        )
        if self.file_path:
            payload["last_file_path"] = self.file_path
        try:
            payload["last_sheet_name"] = self.current_sheet.get()
        except Exception:
            pass
        try:
            with open(SETTINGS_FILE, "w") as f:
                json.dump(payload, f, indent=2)
        except Exception:
            pass
        self.settings = payload

    def _load_window_geom(self):
        try:
            g = self.settings.get("window_geometry", "")
            return g if g else None
        except Exception:
            return None

    def _remember_geom(self, *_):
        if self.state() == "normal":
            self.settings["window_geometry"] = self.geometry()

    def _on_close(self):
        try:
            self._save_settings()
        finally:
            try:
                plt.close("all")
            except Exception:
                pass
            self.destroy()

    # ---------- UI ----------
    def _build_ui(self):
        self.bind("<Configure>", self._remember_geom)

        nb = ttk.Notebook(self)
        nb.pack(fill="both", expand=True, padx=10, pady=10)
        self.nb = nb

        self.tab_data = ttk.Frame(nb)
        self.tab_plot = ttk.Frame(nb)

        nb.add(self.tab_data, text="Data")
        nb.add(self.tab_plot, text="Plot")

        self._build_tab_data()
        self._build_tab_plot()

        # Bottom buttons
        bottom = ttk.Frame(self)
        bottom.pack(fill="x", padx=10, pady=(0, 10))
        ttk.Button(
            bottom, text="Render Plot in New Tab", command=self.render_plots
        ).pack(side="left")
        ttk.Button(bottom, text="Export PNG/SVG/PDF", command=self.export_plots).pack(
            side="left", padx=8
        )
        ttk.Button(bottom, text="Save Settings", command=self._save_settings).pack(
            side="left", padx=8
        )
        ttk.Button(
            bottom, text="Close All Figures", command=self._close_all_figures
        ).pack(side="left", padx=8)
        ttk.Button(bottom, text="Save Preset", command=self._save_preset).pack(
            side="left", padx=8
        )
        ttk.Button(bottom, text="Load Preset", command=self._load_preset).pack(
            side="left", padx=4
        )
        ttk.Button(
            bottom,
            text="Export to PowerPoint",
            command=lambda: self._export_pptx(
                [plt.figure(n) for n in plt.get_fignums()]
            ),
        ).pack(side="left", padx=8)

    def _display_rendered_figures(self, figs: list[plt.Figure]):
        if not figs:
            return
        if self.nb is None:
            try:
                plt.show(block=False)
            except Exception:
                pass
            return

        self.render_tab_counter += 1
        tab_title = f"Render {self.render_tab_counter}"
        tab = ttk.Frame(self.nb)
        tab_info = {"tab": tab, "canvases": [], "toolbars": [], "figures": figs}

        header = ttk.Frame(tab)
        header.pack(fill="x", padx=8, pady=(8, 4))
        ttk.Label(
            header,
            text=f"{len(figs)} figure(s) rendered",
            font=("Times New Roman", 12, "bold"),
        ).pack(side="left")
        ttk.Button(
            header,
            text="Close Tab",
            command=lambda info=tab_info: self._close_render_tab(info),
        ).pack(side="right")

        if len(figs) == 1:
            container = ttk.Frame(tab)
            container.pack(fill="both", expand=True, padx=8, pady=(0, 8))
            canvas = FigureCanvasTkAgg(figs[0], master=container)
            widget = canvas.get_tk_widget()
            widget.pack(fill="both", expand=True)
            toolbar = NavigationToolbar2Tk(canvas, container)
            toolbar.update()
            toolbar.pack(fill="x", pady=(4, 0))
            canvas.draw_idle()
            tab_info["canvases"].append(canvas)
            tab_info["toolbars"].append(toolbar)
        else:
            fig_nb = ttk.Notebook(tab)
            fig_nb.pack(fill="both", expand=True, padx=8, pady=(0, 8))
            tab_info["notebook"] = fig_nb
            for idx, fig in enumerate(figs, start=1):
                fig_frame = ttk.Frame(fig_nb)
                fig_nb.add(fig_frame, text=f"Figure {idx}")
                canvas = FigureCanvasTkAgg(fig, master=fig_frame)
                widget = canvas.get_tk_widget()
                widget.pack(fill="both", expand=True)
                toolbar = NavigationToolbar2Tk(canvas, fig_frame)
                toolbar.update()
                toolbar.pack(fill="x", pady=(4, 0))
                canvas.draw_idle()
                tab_info["canvases"].append(canvas)
                tab_info["toolbars"].append(toolbar)

        self.nb.add(tab, text=tab_title)
        self.nb.select(tab)
        self.render_tabs.append(tab_info)

    def _close_render_tab(self, tab_info: dict):
        if tab_info not in self.render_tabs:
            return
        canvases = tab_info.get("canvases", [])
        for canvas in canvases:
            try:
                fig = canvas.figure
            except Exception:
                fig = None
            try:
                widget = canvas.get_tk_widget()
                widget.destroy()
            except Exception:
                pass
            try:
                canvas._tkcanvas.destroy()
            except Exception:
                pass
            if fig is not None:
                try:
                    plt.close(fig)
                except Exception:
                    pass
        for toolbar in tab_info.get("toolbars", []):
            try:
                toolbar.destroy()
            except Exception:
                pass
        try:
            self.nb.forget(tab_info["tab"])
        except Exception:
            pass
        try:
            tab_info["tab"].destroy()
        except Exception:
            pass
        self.render_tabs.remove(tab_info)

    def _close_all_figures(self):
        for tab_info in list(self.render_tabs):
            self._close_render_tab(tab_info)
        try:
            plt.close("all")
        except Exception:
            pass

    def _build_tab_data(self):
        f = self.tab_data
        for i in range(4):
            f.grid_columnconfigure(i, weight=1)

        ttk.Label(f, text="Data file (CSV or Excel):").grid(
            row=0, column=0, sticky="w", padx=6, pady=6
        )
        self.e_path = ttk.Entry(f)
        self.e_path.grid(row=0, column=1, columnspan=2, sticky="ew", padx=6, pady=6)
        ttk.Button(f, text="Browse…", command=self._browse_file).grid(
            row=0, column=3, sticky="e", padx=6, pady=6
        )

        ttk.Button(
            f,
            text="Load / Refresh",
            command=lambda: self._load_file(self.e_path.get().strip()),
        ).grid(row=1, column=3, sticky="e", padx=6, pady=6)

        ttk.Label(f, text="Sheet (Excel only):").grid(
            row=1, column=0, sticky="w", padx=6, pady=6
        )
        self.cb_sheet = ttk.Combobox(
            f,
            values=self.sheet_names,
            textvariable=self.current_sheet,
            state="readonly",
        )
        self.cb_sheet.grid(row=1, column=1, sticky="ew", padx=6, pady=6)
        self.cb_sheet.bind("<<ComboboxSelected>>", self._on_sheet_selection_change)

        load_sheet_frame = ttk.Frame(f)
        load_sheet_frame.grid(row=1, column=2, sticky="w", padx=6, pady=6)

        ttk.Button(
            load_sheet_frame,
            text="Load Sheet",
            command=self._load_sheet,
        ).pack(side="left")

        self._sheet_indicator_canvas = self._create_sheet_status_indicator(
            load_sheet_frame
        )

        self.lbl_status = ttk.Label(f, text="No data loaded.")
        self.lbl_status.grid(row=2, column=0, columnspan=4, sticky="w", padx=6, pady=10)

    def _create_sheet_status_indicator(self, parent):
        canvas = tk.Canvas(
            parent, width=14, height=14, highlightthickness=0, borderwidth=0
        )
        canvas.pack(side="left", padx=(6, 0))
        self._sheet_indicator_canvas = canvas
        self._update_sheet_indicator(self._sheet_indicator_state)
        return canvas

    def _update_sheet_indicator(self, loaded: bool) -> None:
        self._sheet_indicator_state = bool(loaded)
        canvas = self._sheet_indicator_canvas
        if canvas is None or not canvas.winfo_exists():
            return
        canvas.delete("all")
        radius = 4
        fill_color = "#2da44e" if loaded else "#d73a49"
        canvas.create_oval(
            2, 2, 2 + radius * 2, 2 + radius * 2, fill=fill_color, outline=""
        )

    def _on_sheet_selection_change(self, *_):
        self._update_sheet_indicator(False)

    def _build_tab_plot(self):
        f = self.tab_plot

        # The tab content itself is one big Panedwindow with two panes (left controls + right preview)
        f.grid_rowconfigure(0, weight=1)
        f.grid_columnconfigure(0, weight=1)

        paned = ttk.Panedwindow(f, orient="horizontal")
        paned.grid(row=0, column=0, sticky="nsew", padx=6, pady=6)

        # ----- LEFT PANE (resizable) -----
        left_pane = ttk.Frame(paned)
        paned.add(left_pane)  # ADD ONCE

        # Configure pane options using the correct API (no re-adding)
        try:
            paned.pane(left_pane, weight=1)
        except Exception:
            pass
        try:
            paned.pane(left_pane, minsize=LEFT_COL_MIN)
        except Exception:
            pass

        # Canvas + vertical scrollbar inside left pane
        self.left_canvas = tk.Canvas(left_pane, highlightthickness=0, borderwidth=0)
        vbar = ttk.Scrollbar(
            left_pane, orient="vertical", command=self.left_canvas.yview
        )
        self.left_scrollframe = ttk.Frame(self.left_canvas)
        self.left_canvas_window = self.left_canvas.create_window(
            (0, 0), window=self.left_scrollframe, anchor="nw"
        )
        self.left_canvas.configure(yscrollcommand=vbar.set)

        self.left_canvas.grid(row=0, column=0, sticky="nsew")
        vbar.grid(row=0, column=1, sticky="ns")
        left_pane.grid_rowconfigure(0, weight=1)
        left_pane.grid_columnconfigure(0, weight=1)

        # Ensure canvas scrollregion follows content
        self.left_scrollframe.bind(
            "<Configure>",
            lambda e: self.left_canvas.configure(
                scrollregion=self.left_canvas.bbox("all")
            ),
        )

        # Make the canvas window width follow the current left pane width (minus scrollbar)
        def _sync_inner_width(event=None):
            bar_w = vbar.winfo_width() or SCROLLBAR_FALLBACK_WIDTH
            w = max(120, left_pane.winfo_width() - bar_w)
            self.left_canvas.itemconfigure(self.left_canvas_window, width=w)

        left_pane.bind("<Configure>", _sync_inner_width)
        self.left_canvas.bind(
            "<Configure>", _sync_inner_width
        )  # NEW: catch canvas resizes too
        self.after(50, _sync_inner_width)  # do an initial sync

        # Set initial sash position after the panedwindow is realized
        def _place_sash():
            try:
                paned.sashpos(0, LEFT_COL_INIT)
            except Exception:
                try:
                    paned.sash_place(0, LEFT_COL_INIT, 0)
                except Exception:
                    pass

        self.after(80, _place_sash)

        # Mouse wheel scrolling over the left pane only
        def _on_mousewheel(event):
            self.left_canvas.yview_scroll(int(-1 * (event.delta / 120)), "units")

        def _on_linux_up(event):
            self.left_canvas.yview_scroll(-1, "units")

        def _on_linux_dn(event):
            self.left_canvas.yview_scroll(1, "units")

        self.left_canvas.bind(
            "<Enter>",
            lambda e: (
                self.left_canvas.bind_all("<MouseWheel>", _on_mousewheel),
                self.left_canvas.bind("<Button-4>", _on_linux_up),
                self.left_canvas.bind("<Button-5>", _on_linux_dn),
            ),
        )
        self.left_canvas.bind(
            "<Leave>",
            lambda e: (
                self.left_canvas.unbind_all("<MouseWheel>"),
                self.left_canvas.unbind("<Button-4>"),
                self.left_canvas.unbind("<Button-5>"),
            ),
        )

        # Everything below builds **into** left_scrollframe
        left = self.left_scrollframe

        # ----- RIGHT PANE (preview) -----
        right = ttk.Frame(paned)
        paned.add(right)  # add once
        try:
            paned.pane(right, weight=1)  # correct API
        except Exception:
            pass  # don't re-add

        right.grid_rowconfigure(1, weight=1)
        right.grid_columnconfigure(0, weight=1)

        # ---------- Controls (left) ----------
        lf_type = ttk.Labelframe(left, text="Plot Type")
        lf_type.pack(fill="x", padx=6, pady=(6, 4))
        ttk.Label(lf_type, text="Choose a layout:").grid(
            row=0, column=0, sticky="w", padx=6, pady=4
        )
        self.cb_plot_type = ttk.Combobox(
            lf_type,
            state="readonly",
            width=28,
            textvariable=self.plot_type,
            values=[
                "single",
                "single_right",
                "grid",
                "grid_right",
                "facet",
                "facet_right",
            ],
        )
        self.cb_plot_type.grid(row=0, column=1, sticky="w", padx=6, pady=4)
        lf_type.grid_columnconfigure(1, weight=1)

        # Wrap text adjusts as left pane changes size
        self.lbl_type_desc = ttk.Label(
            lf_type,
            foreground="#333",
            text=self._plot_type_description(self.plot_type.get()),
            wraplength=LEFT_COL_INIT - 60,
            justify="left",
        )
        self.lbl_type_desc.grid(
            row=1, column=0, columnspan=2, sticky="w", padx=6, pady=(0, 6)
        )
        self.lbl_type_tip = ttk.Label(
            lf_type,
            foreground="#444",
            text=self._plot_type_tip(self.plot_type.get()),
            wraplength=LEFT_COL_INIT - 60,
            justify="left",
        )
        self.lbl_type_tip.grid(
            row=2, column=0, columnspan=2, sticky="w", padx=6, pady=(0, 6)
        )

        # update wraplength with pane width
        def _update_wrap(*_):
            bar_w = vbar.winfo_width() or SCROLLBAR_FALLBACK_WIDTH
            w = max(200, left_pane.winfo_width() - bar_w - 30)
            self.lbl_type_desc.configure(wraplength=w)
            self.lbl_type_tip.configure(wraplength=w)

        left_pane.bind("<Configure>", _update_wrap)

        lf_cols = ttk.Labelframe(left, text="Columns")
        lf_cols.pack(fill="x", padx=6, pady=4)
        ttk.Label(lf_cols, text="X Column (Required)").grid(
            row=0, column=0, sticky="w", padx=6, pady=4
        )
        self.cb_x = ttk.Combobox(lf_cols, values=self.columns, state="readonly")
        self.cb_x.grid(row=0, column=1, sticky="ew", padx=6, pady=4)
        lf_cols.grid_columnconfigure(1, weight=1)
        ttk.Label(lf_cols, text="Y Columns (one or more)").grid(
            row=1, column=0, sticky="nw", padx=6, pady=4
        )
        self.lb_y = tk.Listbox(
            lf_cols, selectmode="extended", exportselection=False, height=8
        )
        self.lb_y.grid(row=1, column=1, sticky="nsew", padx=6, pady=4)
        lf_cols.grid_rowconfigure(1, weight=1)
        btn_row = ttk.Frame(lf_cols)
        btn_row.grid(row=2, column=1, sticky="w", padx=6, pady=(2, 6))
        ttk.Button(
            btn_row, text="Edit series styles…", command=self._open_style_manager
        ).pack(side="left")
        ttk.Button(btn_row, text="Filter…", command=self._open_filter_dialog).pack(
            side="left", padx=6
        )
        ttk.Label(lf_cols, text="Error bar (yerr) column").grid(
            row=3, column=0, sticky="w", padx=6, pady=4
        )
        self.cb_yerr = ttk.Combobox(
            lf_cols,
            values=["None"] + self.columns,
            textvariable=self.yerr_col,
            state="readonly",
            width=18,
        )
        self.cb_yerr.grid(row=3, column=1, sticky="w", padx=6, pady=4)

        self.lf_right = ttk.Labelframe(
            left, text="Secondary Right Y (for *_right types)"
        )
        self.lf_right.pack(fill="x", padx=6, pady=4)
        ttk.Label(
            self.lf_right, text="Choose ONE series to plot on the right axis:"
        ).pack(side="left", padx=6, pady=6)
        self.cb_right = ttk.Combobox(
            self.lf_right,
            values=["None"],
            textvariable=self.right_axis_series,
            state="readonly",
            width=28,
        )
        self.cb_right.pack(side="left", padx=6, pady=6)

        self.lf_facet = ttk.Labelframe(
            left, text="Facet (split data by a category column)"
        )
        self.lf_facet.pack(fill="x", padx=6, pady=4)
        ttk.Label(self.lf_facet, text="Facet by (category column)").grid(
            row=0, column=0, sticky="w", padx=6, pady=4
        )
        self.cb_facet = ttk.Combobox(
            self.lf_facet, values=["None"], state="readonly", width=28
        )
        self.cb_facet.grid(row=0, column=1, sticky="ew", padx=6, pady=4)
        self.lf_facet.grid_columnconfigure(1, weight=1)
        ttk.Label(
            self.lf_facet,
            text="Tip: Faceting makes one subplot per unique category value.",
        ).grid(row=1, column=0, columnspan=2, sticky="w", padx=6, pady=(0, 6))

        lf_titles = ttk.Labelframe(left, text="Titles")
        lf_titles.pack(fill="x", padx=6, pady=4)
        ttk.Label(lf_titles, text="Title").grid(
            row=0, column=0, sticky="w", padx=6, pady=4
        )
        ttk.Entry(lf_titles, textvariable=self.title_text).grid(
            row=0, column=1, sticky="ew", padx=6, pady=4
        )
        ttk.Label(lf_titles, text="Suptitle").grid(
            row=1, column=0, sticky="w", padx=6, pady=4
        )
        ttk.Entry(lf_titles, textvariable=self.suptitle_text).grid(
            row=1, column=1, sticky="ew", padx=6, pady=4
        )
        lf_titles.grid_columnconfigure(1, weight=1)

        # --- Style & Layout (boxed subsections) ---
        lf_style = ttk.Labelframe(left, text="Style & Layout")
        lf_style.pack(fill="x", padx=6, pady=4)

        # Optional: a little padding and bold labels for inner boxes
        try:
            style = ttk.Style(self)
            style.configure("Inner.TLabelframe", padding=(8, 6, 8, 6))
            style.configure(
                "Inner.TLabelframe.Label", font=("TkDefaultFont", 9, "bold")
            )
        except Exception:
            pass

        # Make two columns for the top row of boxes
        lf_style.grid_columnconfigure(0, weight=1)
        lf_style.grid_columnconfigure(1, weight=1)

        # ---------- Box 1 (left): Plot appearance ----------
        lf_appear = ttk.Labelframe(
            lf_style, text="Plot appearance", style="Inner.TLabelframe"
        )
        lf_appear.grid(row=0, column=0, sticky="nsew", padx=(6, 3), pady=(6, 6))
        lf_appear.grid_columnconfigure(1, weight=1)

        ttk.Label(lf_appear, text="Kind").grid(
            row=0, column=0, sticky="e", padx=(0, 6), pady=2
        )
        ttk.Combobox(
            lf_appear,
            values=["line", "scatter", "bar", "area", "step"],
            textvariable=self.plot_kind,
            state="readonly",
            width=12,
        ).grid(row=0, column=1, sticky="w", pady=2)

        ttk.Label(lf_appear, text="Linewidth").grid(
            row=1, column=0, sticky="e", padx=(0, 6), pady=2
        )
        ttk.Entry(lf_appear, textvariable=self.linewidth, width=10).grid(
            row=1, column=1, sticky="w", pady=2
        )

        ttk.Label(lf_appear, text="Markersize").grid(
            row=2, column=0, sticky="e", padx=(0, 6), pady=2
        )
        ttk.Entry(lf_appear, textvariable=self.markersize, width=10).grid(
            row=2, column=1, sticky="w", pady=2
        )

        ttk.Label(lf_appear, text="Theme").grid(
            row=3, column=0, sticky="e", padx=(0, 6), pady=2
        )
        ttk.Combobox(
            lf_appear,
            values=["classic", "minimal", "dark", "journal"],
            textvariable=self.theme,
            state="readonly",
            width=14,
        ).grid(row=3, column=1, sticky="w", pady=2)

        # ---------- Box 2 (right): Legend & grid ----------
        lf_layout = ttk.Labelframe(
            lf_style, text="Legend & grid", style="Inner.TLabelframe"
        )
        lf_layout.grid(row=0, column=1, sticky="nsew", padx=(3, 6), pady=(6, 6))
        lf_layout.grid_columnconfigure(1, weight=1)

        ttk.Label(lf_layout, text="Legend location").grid(
            row=0, column=0, sticky="e", padx=(0, 6), pady=2
        )
        ttk.Combobox(
            lf_layout,
            values=[
                "best",
                "upper right",
                "upper left",
                "lower left",
                "lower right",
                "right",
                "center left",
                "center right",
                "lower center",
                "upper center",
                "center",
                "outside bottom",
            ],
            textvariable=self.legend_loc,
            state="readonly",
            width=16,
        ).grid(row=0, column=1, sticky="w", pady=2)

        ttk.Label(lf_layout, text="Grid columns").grid(
            row=1, column=0, sticky="e", padx=(0, 6), pady=2
        )
        ttk.Spinbox(lf_layout, from_=1, to=6, textvariable=self.ncols, width=8).grid(
            row=1, column=1, sticky="w", pady=2
        )

        ttk.Checkbutton(
            lf_layout, text="Decimate preview for speed", variable=self.decimate_preview
        ).grid(row=2, column=0, columnspan=2, sticky="w", pady=(6, 0))

        # NEW: grid toggle
        ttk.Checkbutton(
            lf_layout,
            text="Show grid lines",
            variable=self.show_grid,
            command=self._schedule_preview,
        ).grid(row=3, column=0, columnspan=2, sticky="w", pady=(4, 0))

        # ---------- Box 3 (full width): Axis label formatting ----------
        lf_fmt = ttk.Labelframe(
            lf_style, text="Axis label formatting", style="Inner.TLabelframe"
        )
        lf_fmt.grid(row=1, column=0, columnspan=2, sticky="ew", padx=6, pady=(0, 6))
        lf_fmt.grid_columnconfigure(1, weight=1)
        lf_fmt.grid_columnconfigure(3, weight=1)

        # === NEW: Axis label text inputs ===
        ttk.Label(lf_fmt, text="X axis label").grid(
            row=0, column=0, sticky="e", padx=6, pady=4
        )
        ttk.Entry(lf_fmt, textvariable=self.x_label).grid(
            row=0, column=1, sticky="ew", padx=6, pady=4
        )

        ttk.Label(lf_fmt, text="Left Y label").grid(
            row=1, column=0, sticky="e", padx=6, pady=4
        )
        ttk.Entry(lf_fmt, textvariable=self.y_left_label).grid(
            row=1, column=1, sticky="ew", padx=6, pady=4
        )

        # Right-Y label (only shown for *_right plot types)
        self.frm_y_right_label = ttk.Frame(lf_fmt)
        self.frm_y_right_label.grid(
            row=2, column=0, columnspan=2, sticky="ew", padx=0, pady=(0, 2)
        )
        ttk.Label(self.frm_y_right_label, text="Right Y label").grid(
            row=0, column=0, sticky="e", padx=6, pady=4
        )
        self.e_y_right_label = ttk.Entry(
            self.frm_y_right_label, textvariable=self.y_right_label
        )
        self.e_y_right_label.grid(row=0, column=1, sticky="ew", padx=6, pady=4)

        # Let the entry column stretch
        lf_fmt.grid_columnconfigure(1, weight=1)

        ttk.Label(lf_fmt, text="Left Y format").grid(
            row=3, column=0, sticky="e", padx=(0, 6), pady=2
        )
        ttk.Combobox(
            lf_fmt,
            values=["plain", "sci", "percent", "thousands"],
            textvariable=self.y_left_format,
            state="readonly",
            width=14,
        ).grid(row=3, column=1, sticky="w", pady=2)

        ttk.Label(lf_fmt, text="Right Y format").grid(
            row=3, column=2, sticky="e", padx=(18, 6), pady=2
        )
        ttk.Combobox(
            lf_fmt,
            values=["plain", "sci", "percent", "thousands"],
            textvariable=self.y_right_format,
            state="readonly",
            width=14,
        ).grid(row=3, column=3, sticky="w", pady=2)

        lf_ranges = ttk.Labelframe(left, text="Axis Ranges")
        lf_ranges.pack(fill="x", padx=6, pady=4)
        self.log_x = tk.BooleanVar(value=self.settings.get("log_x", False))
        self.log_y_left = tk.BooleanVar(value=self.settings.get("log_y_left", False))
        self.log_y_right = tk.BooleanVar(value=self.settings.get("log_y_right", False))
        ttk.Checkbutton(lf_ranges, text="Log X", variable=self.log_x).grid(
            row=3, column=0, sticky="w", padx=6, pady=4
        )
        ttk.Checkbutton(lf_ranges, text="Log Left Y", variable=self.log_y_left).grid(
            row=3, column=1, sticky="w", padx=6, pady=4
        )
        ttk.Checkbutton(lf_ranges, text="Log Right Y", variable=self.log_y_right).grid(
            row=3, column=2, sticky="w", padx=6, pady=4
        )
        self.x_tick_rotation = tk.IntVar(
            value=int(self.settings.get("x_tick_rotation", 0))
        )
        ttk.Label(lf_ranges, text="X tick rotation").grid(
            row=3, column=3, sticky="e", padx=6, pady=4
        )
        ttk.Spinbox(
            lf_ranges, from_=-90, to=90, textvariable=self.x_tick_rotation, width=6
        ).grid(row=3, column=4, sticky="w", padx=6, pady=4)
        ttk.Checkbutton(
            lf_ranges,
            text="Auto X",
            variable=self.auto_x,
            command=self._sync_range_states,
        ).grid(row=0, column=0, sticky="w", padx=6, pady=4)
        ttk.Label(lf_ranges, text="X min").grid(
            row=0, column=1, sticky="e", padx=6, pady=4
        )
        self.e_xmin = ttk.Entry(lf_ranges, textvariable=self.x_min, width=10)
        self.e_xmin.grid(row=0, column=2, sticky="w", padx=6, pady=4)
        ttk.Label(lf_ranges, text="X max").grid(
            row=0, column=3, sticky="e", padx=6, pady=4
        )
        self.e_xmax = ttk.Entry(lf_ranges, textvariable=self.x_max, width=10)
        self.e_xmax.grid(row=0, column=4, sticky="w", padx=6, pady=4)
        ttk.Checkbutton(
            lf_ranges,
            text="Auto Left Y",
            variable=self.auto_y_left,
            command=self._sync_range_states,
        ).grid(row=1, column=0, sticky="w", padx=6, pady=4)
        ttk.Label(lf_ranges, text="Yₗ min").grid(
            row=1, column=1, sticky="e", padx=6, pady=4
        )
        self.e_ylmin = ttk.Entry(lf_ranges, textvariable=self.y_left_min, width=10)
        self.e_ylmin.grid(row=1, column=2, sticky="w", padx=6, pady=4)
        ttk.Label(lf_ranges, text="Yₗ max").grid(
            row=1, column=3, sticky="e", padx=6, pady=4
        )
        self.e_ylmax = ttk.Entry(lf_ranges, textvariable=self.y_left_max, width=10)
        self.e_ylmax.grid(row=1, column=4, sticky="w", padx=6, pady=4)
        ttk.Checkbutton(
            lf_ranges,
            text="Auto Right Y",
            variable=self.auto_y_right,
            command=self._sync_range_states,
        ).grid(row=2, column=0, sticky="w", padx=6, pady=4)
        ttk.Label(lf_ranges, text="Yᵣ min").grid(
            row=2, column=1, sticky="e", padx=6, pady=4
        )
        self.e_yrmin = ttk.Entry(lf_ranges, textvariable=self.y_right_min, width=10)
        self.e_yrmin.grid(row=2, column=2, sticky="w", padx=6, pady=4)
        ttk.Label(lf_ranges, text="Yᵣ max").grid(
            row=2, column=3, sticky="e", padx=6, pady=4
        )
        self.e_yrmax = ttk.Entry(lf_ranges, textvariable=self.y_right_max, width=10)
        self.e_yrmax.grid(row=2, column=4, sticky="w", padx=6, pady=4)

        # Keep the manual range boxes in sync with the Auto toggles
        self.auto_x.trace_add("write", lambda *_: self._sync_range_states())
        self.auto_y_left.trace_add("write", lambda *_: self._sync_range_states())
        self.auto_y_right.trace_add("write", lambda *_: self._sync_range_states())

        lf_ticks = ttk.Labelframe(left, text="Ticks")
        lf_ticks.pack(fill="x", padx=6, pady=(4, 8))
        ttk.Checkbutton(
            lf_ticks, text="Auto X ticks", variable=self.auto_time_ticks
        ).grid(row=0, column=0, sticky="w", padx=6, pady=4)
        ttk.Label(lf_ticks, text="X major").grid(
            row=0, column=1, sticky="e", padx=6, pady=4
        )
        ttk.Entry(lf_ticks, textvariable=self.xmaj, width=8).grid(
            row=0, column=2, sticky="w", padx=6, pady=4
        )
        ttk.Label(lf_ticks, text="X minor").grid(
            row=0, column=3, sticky="e", padx=6, pady=4
        )
        ttk.Entry(lf_ticks, textvariable=self.xminr, width=8).grid(
            row=0, column=4, sticky="w", padx=6, pady=4
        )
        ttk.Checkbutton(
            lf_ticks, text="Auto Left Y ticks", variable=self.auto_y_ticks
        ).grid(row=1, column=0, sticky="w", padx=6, pady=4)
        ttk.Label(lf_ticks, text="Y major").grid(
            row=1, column=1, sticky="e", padx=6, pady=4
        )
        ttk.Entry(lf_ticks, textvariable=self.ymaj, width=8).grid(
            row=1, column=2, sticky="w", padx=6, pady=4
        )
        ttk.Label(lf_ticks, text="Y minor").grid(
            row=1, column=3, sticky="e", padx=6, pady=4
        )
        ttk.Entry(lf_ticks, textvariable=self.yminr, width=8).grid(
            row=1, column=4, sticky="w", padx=6, pady=4
        )

        lf_refs = ttk.Labelframe(left, text="References (lines & bands)")
        lf_refs.pack(fill="x", padx=6, pady=(4, 8))
        rbar = ttk.Frame(lf_refs)
        rbar.pack(anchor="w", padx=6, pady=4)
        ttk.Button(
            rbar, text="Add H-Line", command=lambda: self._add_reference("hline")
        ).pack(side="left")
        ttk.Button(
            rbar, text="Add V-Line", command=lambda: self._add_reference("vline")
        ).pack(side="left", padx=5)
        ttk.Button(
            rbar, text="Add Y-Band", command=lambda: self._add_reference("band_y")
        ).pack(side="left", padx=5)
        ttk.Button(
            rbar, text="Add X-Band", command=lambda: self._add_reference("band_x")
        ).pack(side="left", padx=5)
        ttk.Button(rbar, text="Clear", command=self._clear_references).pack(
            side="left", padx=10
        )
        ttk.Label(lf_refs, text="Facets: limit to top-K categories (0 = all)").pack(
            anchor="w", padx=6
        )
        ttk.Spinbox(
            lf_refs, from_=0, to=999, textvariable=self.top_k_facets, width=6
        ).pack(anchor="w", padx=6, pady=(0, 4))

        # ---------- Preview (right) ----------
        ttk.Label(
            right, text="Live Preview", font=("Times New Roman", 13, "bold")
        ).grid(row=0, column=0, sticky="w", padx=6, pady=(0, 0))
        self.preview_fig = plt.Figure(figsize=(6.6, 4.8))
        self.preview_canvas = FigureCanvasTkAgg(self.preview_fig, master=right)
        self.preview_widget = self.preview_canvas.get_tk_widget()
        self.preview_widget.grid(row=1, column=0, sticky="nsew", padx=6, pady=6)
        ttk.Button(right, text="Refresh Preview", command=self.draw_preview).grid(
            row=2, column=0, sticky="e", padx=6, pady=(0, 6)
        )

        stats_frame = ttk.Labelframe(right, text="Stats (selected data)")
        stats_frame.grid(row=4, column=0, sticky="nsew", padx=6, pady=(0, 6))
        right.grid_rowconfigure(4, weight=0)
        self.stats_text = tk.Text(stats_frame, height=6, wrap="word")
        self.stats_text.pack(fill="both", expand=True, padx=6, pady=6)

        toolbar_frame = ttk.Frame(right)
        toolbar_frame.grid(row=3, column=0, sticky="ew", padx=6, pady=(0, 6))
        self.preview_toolbar = NavigationToolbar2Tk(self.preview_canvas, toolbar_frame)
        self.preview_toolbar.update()

        self.cb_plot_type.bind("<<ComboboxSelected>>", self._on_plot_type_change)
        self.plot_type.trace_add("write", lambda *args: self._on_plot_type_change())
        self._wire_live_preview(left)
        self._on_plot_type_change()
        self._sync_range_states()

    # ---- Helpers for plot type description & contextual UI
    def _plot_type_description(self, key: str) -> str:
        d = {
            "single": "Single axes: plots all selected Y series against X on the same (left) Y-axis.",
            "single_right": "Single axes + Right Y: same as Single, plus ONE selected series on a secondary (right) Y-axis.",
            "grid": "Subplot grid: one subplot per Y series.",
            "grid_right": "Subplot grid + Right Y: each subplot shows its Y series on the left axis; the chosen right-axis series overlays on the right axis in each subplot (if present).",
            "facet": "Facet grid: choose a category column; creates one subplot per unique category. All selected Y series are drawn in each facet.",
            "facet_right": "Facet grid + Right Y: same as Facet, with ONE series plotted on a secondary right axis in each facet.",
        }
        return d.get(key, "")

    def _plot_type_tip(self, key: str) -> str:
        tips = {
            "single": "Explainer: All selected Y series are drawn together on ONE set of axes. Use this when you want to compare series directly on the same scale.",
            "single_right": "Explainer: Same as Single, but ONE chosen series is drawn on a secondary (right) Y-axis. Use when one series has a very different scale.",
            "grid": "Explainer: Builds a small-multiples grid with ONE subplot per selected Y series. Use this to compare shapes/patterns without overlap.",
            "grid_right": "Explainer: Like Grid, but overlays ONE chosen series on a right Y-axis in each subplot. Useful if a common reference series should appear with different scales.",
            "facet": "Explainer: Pick a category column to facet on. The data is split by category and each category gets its own subplot. All selected Y series appear in each facet.",
            "facet_right": "Explainer: Like Facet, plus ONE chosen series on a secondary (right) Y-axis in every facet. Helps compare categories while keeping a reference series readable.",
        }
        return tips.get(key, "")

    def _on_plot_type_change(self, *_):
        key = self.plot_type.get()
        self.lbl_type_desc.configure(text=self._plot_type_description(key))
        self.lbl_type_tip.configure(text=self._plot_type_tip(key))
        right_needed = key.endswith("_right")
        facet_needed = key.startswith("facet")
        self.lf_right.pack_forget()
        self.lf_facet.pack_forget()
        if right_needed:
            self.lf_right.pack(fill="x", padx=6, pady=4)
        if facet_needed:
            self.lf_facet.pack(fill="x", padx=6, pady=4)
        state = "normal" if right_needed else "disabled"
        for w in self.lf_right.winfo_children():
            try:
                w.configure(state=state)
            except Exception:
                pass
        # Show/hide the Right Y label row inside Axis label formatting
        try:
            if right_needed:
                self.frm_y_right_label.grid()  # ensure it's visible
            else:
                self.frm_y_right_label.grid_remove()
        except Exception:
            pass

        self._schedule_preview()

    def _wire_live_preview(self, root):
        def schedule(*_):
            self._schedule_preview()

        self._bind_recursive(root, schedule)
        self.lb_y.bind("<<ListboxSelect>>", schedule)

    def _bind_recursive(self, widget, callback):
        if isinstance(widget, ttk.Combobox):
            widget.bind("<<ComboboxSelected>>", callback)
        cls = widget.winfo_class().lower()
        if "entry" in cls:
            widget.bind("<FocusOut>", callback)
            widget.bind("<Return>", callback)
        if isinstance(widget, ttk.Spinbox):
            widget.bind("<FocusOut>", callback)
            widget.bind("<Return>", callback)
        # ↓ Replace your existing check/radio block with this:
        if isinstance(widget, (ttk.Checkbutton, ttk.Radiobutton)):
            try:
                # Only set a command if one is not already present
                if not widget.cget("command"):
                    widget.configure(command=callback)
            except Exception:
                pass
        for child in widget.winfo_children():
            self._bind_recursive(child, callback)

    def _schedule_preview(self):
        if self._preview_after_id:
            self.after_cancel(self._preview_after_id)
        self._preview_after_id = self.after(
            self._preview_debounce_ms, self.draw_preview
        )

    # ---------- Small dialogs / actions ----------
    def _open_filter_dialog(self):
        top = tk.Toplevel(self)
        top.title("Filter Data")
        top.grab_set()
        ttk.Label(top, text="Pandas query (e.g. Category=='A' and value>10):").pack(
            anchor="w", padx=10, pady=(10, 4)
        )
        e = ttk.Entry(top, width=60)
        e.insert(0, self.filter_query.get())
        e.pack(fill="x", padx=10, pady=4)

        def apply_and_close():
            self.filter_query.set(e.get().strip())
            self._schedule_preview()
            top.destroy()

        ttk.Button(top, text="Apply", command=apply_and_close).pack(pady=8)

    def _open_style_manager(self):
        ycols = self._get_selected_y()
        if not ycols:
            messagebox.showinfo("Series styles", "Select at least one Y column first.")
            return
        top = tk.Toplevel(self)
        top.title("Series Styles")
        top.grab_set()
        frm = ttk.Frame(top)
        frm.pack(fill="both", expand=True, padx=10, pady=10)
        headers = [
            "Series",
            "Visible",
            "Color",
            "Marker",
            "LW",
            "MS",
            "Scale",
            "Offset",
            "Rolling",
        ]
        for j, h in enumerate(headers):
            ttk.Label(frm, text=h, font=("TkDefaultFont", 9, "bold")).grid(
                row=0, column=j, padx=4, pady=2
            )
        self._style_widgets = []
        for i, col in enumerate(ycols, start=1):
            st = self.series_style.setdefault(col, {})
            ttk.Label(frm, text=col).grid(row=i, column=0, sticky="w", padx=4)
            vis = tk.BooleanVar(value=bool(st.get("visible", True)))
            ttk.Checkbutton(frm, variable=vis).grid(row=i, column=1)

            def pick(c=col):
                init = self.series_style[c].get("color") or "#1f77b4"
                self.series_style[c]["color"] = _safe_color_dialog(initial=init)

            ttk.Button(frm, text="Pick", command=pick).grid(row=i, column=2)
            mvar = tk.StringVar(value=st.get("marker", MARKER_DEFAULT))
            mc = ttk.Combobox(
                frm, values=MARKER_CHOICES, textvariable=mvar, width=4, state="readonly"
            )
            mc.grid(row=i, column=3)
            lw = tk.DoubleVar(value=float(st.get("linewidth", self.linewidth.get())))
            ttk.Spinbox(
                frm, from_=0.1, to=10, increment=0.1, textvariable=lw, width=5
            ).grid(row=i, column=4)
            ms = tk.DoubleVar(value=float(st.get("markersize", self.markersize.get())))
            ttk.Spinbox(
                frm, from_=0.1, to=20, increment=0.5, textvariable=ms, width=5
            ).grid(row=i, column=5)
            sc = tk.DoubleVar(value=float(st.get("scale", 1.0)))
            of = tk.DoubleVar(value=float(st.get("offset", 0.0)))
            ro = tk.IntVar(value=int(st.get("rolling", 0)))
            ttk.Entry(frm, textvariable=sc, width=6).grid(row=i, column=6)
            ttk.Entry(frm, textvariable=of, width=6).grid(row=i, column=7)
            ttk.Entry(frm, textvariable=ro, width=6).grid(row=i, column=8)
            self._style_widgets.append((col, vis, mvar, lw, ms, sc, of, ro))

        def apply_close():
            for col, vis, mvar, lw, ms, sc, of, ro in self._style_widgets:
                st = self.series_style.setdefault(col, {})
                st["visible"] = vis.get()
                st["marker"] = mvar.get()
                st["linewidth"] = lw.get()
                st["markersize"] = ms.get()
                st["scale"] = sc.get()
                st["offset"] = of.get()
                st["rolling"] = ro.get()
            self._save_settings()
            self._schedule_preview()
            top.destroy()

        ttk.Button(top, text="Apply", command=apply_close).pack(pady=10)

    def _add_reference(self, kind):
        d = {}
        if kind == "hline":
            y = simpledialog.askstring("Add horizontal line", "Y value:")
            d = {"type": "hline", "y": y, "alpha": 0.5} if y else {}
        elif kind == "vline":
            x = simpledialog.askstring("Add vertical line", "X value:")
            d = {"type": "vline", "x": x, "alpha": 0.5} if x else {}
        elif kind == "band_y":
            y0 = simpledialog.askstring("Add Y band", "Y min:")
            y1 = simpledialog.askstring("Add Y band", "Y max:")
            d = (
                {"type": "band_y", "y0": y0, "y1": y1, "alpha": 0.15}
                if (y0 and y1)
                else {}
            )
        elif kind == "band_x":
            x0 = simpledialog.askstring("Add X band", "X min:")
            x1 = simpledialog.askstring("Add X band", "X max:")
            d = (
                {"type": "band_x", "x0": x0, "x1": x1, "alpha": 0.15}
                if (x0 and x1)
                else {}
            )
        if d:
            d["color"] = _safe_color_dialog()
            d["label"] = simpledialog.askstring("Label (optional)", "Legend label:")
            self.references.append(d)
            self._schedule_preview()

    def _clear_references(self):
        if messagebox.askyesno("Clear references", "Remove all reference lines/bands?"):
            self.references.clear()
            self._schedule_preview()

    def _save_preset(self):
        name = simpledialog.askstring("Save preset", "Preset name:")
        if not name:
            return
        os.makedirs("presets", exist_ok=True)
        path = os.path.join("presets", f"{name}.json")
        self._save_settings()
        try:
            with open(path, "w") as f:
                json.dump(self.settings, f, indent=2)
            messagebox.showinfo("Preset saved", f"Saved to {path}")
        except Exception as e:
            messagebox.showerror("Preset error", str(e))

    def _load_preset(self):
        path = filedialog.askopenfilename(
            title="Load preset",
            initialdir="presets",
            filetypes=[("Preset JSON", "*.json"), ("All files", "*.*")],
        )
        if not path:
            return
        try:
            with open(path, "r") as f:
                loaded = json.load(f)
            self.settings.update(loaded)
            self.title_text.set(self.settings.get("title_text", ""))
            self.suptitle_text.set(self.settings.get("suptitle_text", ""))
            self.plot_type.set(self.settings.get("plot_type", "single"))
            self.legend_loc.set(self.settings.get("legend_loc", "best"))
            self.plot_kind.set(self.settings.get("plot_kind", "line"))
            self.series_style = self.settings.get("series_style", {})
            self.references = self.settings.get("references", [])
            self.filter_query.set(self.settings.get("filter_query", ""))
            self.decimate_preview.set(self.settings.get("decimate_preview", True))
            self.yerr_col.set(self.settings.get("yerr_col", "None"))
            self.y_left_format.set(self.settings.get("y_left_format", "plain"))
            self.y_right_format.set(self.settings.get("y_right_format", "plain"))
            self.theme.set(self.settings.get("theme", "classic"))
            self.top_k_facets.set(int(self.settings.get("top_k_facets", 0)))
            self._on_plot_type_change()
            self._sync_range_states()
            self._schedule_preview()
            messagebox.showinfo("Preset loaded", os.path.basename(path))
        except Exception as e:
            messagebox.showerror("Preset error", str(e))

    def _export_pptx(self, figs):
        if Presentation is None:
            messagebox.showerror(
                "Export", "python-pptx not installed. Try: pip install python-pptx"
            )
            return
        prs = Presentation()
        blank = prs.slide_layouts[6]
        for idx, fig in enumerate(figs, start=1):
            tmp = os.path.join(os.getcwd(), f"_tmp_{idx}.png")
            fig.savefig(tmp, dpi=EXPORT_DPI, bbox_inches="tight")
            slide = prs.slides.add_slide(blank)
            left, top = Inches(0.5), Inches(0.5)
            slide.shapes.add_picture(tmp, left, top, width=Inches(12.5))
            try:
                os.remove(tmp)
            except Exception:
                pass
        out = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint", "*.pptx")],
            title="Save PowerPoint",
        )
        if not out:
            return
        prs.save(out)
        messagebox.showinfo("Export", f"Saved: {out}")

    # ---------- Data loading ----------
    def _browse_file(self):
        path = filedialog.askopenfilename(
            title="Select Data File",
            filetypes=[("Data files", "*.xlsx *.xls *.csv"), ("All files", "*.*")],
        )
        if not path:
            return
        self.e_path.delete(0, tk.END)
        self.e_path.insert(0, path)

    def _load_file(self, path: str, auto=False):
        if not path or not os.path.exists(path):
            if not auto:
                messagebox.showerror("Missing file", "Select a valid CSV/Excel file.")
            return
        self.file_path = path
        self.e_path.delete(0, tk.END)
        self.e_path.insert(0, path)
        self.sheet_names = []
        self._update_sheet_indicator(False)
        if path.lower().endswith((".xlsx", ".xls")):
            try:
                xls = pd.ExcelFile(path, engine="openpyxl")
                self.sheet_names = xls.sheet_names
            except Exception as e:
                messagebox.showerror("Excel error", f"Could not read workbook: {e}")
                self._update_sheet_indicator(False)
                return
            self.cb_sheet.configure(values=self.sheet_names)
            last_sheet = self.settings.get("last_sheet_name", "")
            if last_sheet in self.sheet_names:
                self.current_sheet.set(last_sheet)
            elif self.sheet_names:
                self.current_sheet.set(self.sheet_names[0])
            self._load_sheet()
        else:
            try:
                self.df = pd.read_csv(path)
            except Exception as e:
                messagebox.showerror("CSV error", f"Could not read CSV: {e}")
                self._update_sheet_indicator(False)
                return
            self._post_load_dataframe()

        self.settings["last_file_path"] = path
        self._save_settings()
        self._schedule_preview()

    def _load_sheet(self):
        if not self.file_path or not self.current_sheet.get():
            messagebox.showerror("Missing info", "Pick a file and a sheet.")
            return
        self._update_sheet_indicator(False)
        try:
            self.df = pd.read_excel(
                self.file_path, sheet_name=self.current_sheet.get(), engine="openpyxl"
            )
        except Exception as e:
            messagebox.showerror("Load error", f"Could not load sheet: {e}")
            self._update_sheet_indicator(False)
            return
        self._post_load_dataframe()
        self._schedule_preview()

    def _post_load_dataframe(self):
        self.columns = list(self.df.columns)
        self.cb_x.configure(values=self.columns)
        self.lb_y.delete(0, tk.END)
        for c in self.columns:
            self.lb_y.insert(tk.END, c)
        self.cb_facet.configure(values=["None"] + self.columns)
        if self.facet_col_saved in (["None"] + self.columns):
            self.cb_facet.set(self.facet_col_saved)
        else:
            self.cb_facet.set("None")
        prev_x = self.settings.get("x_col", "")
        if prev_x in self.columns:
            self.cb_x.set(prev_x)
        elif self.columns:
            self.cb_x.set(self.columns[0])
        prev_y = self.settings.get("y_cols", [])
        for idx, col in enumerate(self.columns):
            if col in prev_y:
                self.lb_y.selection_set(idx)
        self.cb_right.configure(values=["None"] + self.columns)
        if self.right_axis_series.get() not in (["None"] + self.columns):
            self.right_axis_series.set("None")

        self.lbl_status.configure(
            text=f"Loaded: {os.path.basename(self.file_path)}"
            + (
                f" | Sheet: {self.current_sheet.get()}"
                if self.current_sheet.get()
                else ""
            )
            + f" | Rows: {len(self.df)} | Cols: {len(self.columns)}"
        )
        self._update_sheet_indicator(True)

    def _get_selected_y(self):
        return [self.lb_y.get(i) for i in self.lb_y.curselection()]

    # ---------- Plotting helpers ----------
    def _apply_ticks(self, ax, set_x=True):
        if set_x:
            if self.auto_time_ticks.get():
                ax.xaxis.set_major_locator(AutoLocator())
                ax.xaxis.set_minor_locator(AutoMinorLocator())
            else:
                xm = _safe_float(self.xmaj.get())
                xmn = _safe_float(self.xminr.get())
                if xm:
                    ax.xaxis.set_major_locator(MultipleLocator(xm))
                if xmn:
                    ax.xaxis.set_minor_locator(MultipleLocator(xmn))
        if self.auto_y_ticks.get():
            ax.yaxis.set_major_locator(AutoLocator())
            ax.yaxis.set_minor_locator(AutoMinorLocator())
        else:
            ym = _safe_float(self.ymaj.get())
            ymn = _safe_float(self.yminr.get())
            if ym:
                ax.yaxis.set_major_locator(MultipleLocator(ym))
            if ymn:
                ax.yaxis.set_minor_locator(MultipleLocator(ymn))
        ax.minorticks_on()
        ax.tick_params(axis="both", which="major", labelsize=TICK_LABELSIZE)

    def _apply_limits(self, ax, which="left", set_x=True):
        if set_x and not self.auto_x.get():
            xmin = _safe_float(self.x_min.get())
            xmax = _safe_float(self.x_max.get())
            if xmin is not None and xmax is not None:
                ax.set_xlim(xmin, xmax)
        if which == "left":
            if not self.auto_y_left.get():
                ymin = _safe_float(self.y_left_min.get())
                ymax = _safe_float(self.y_left_max.get())
                if ymin is not None and ymax is not None:
                    ax.set_ylim(ymin, ymax)
        else:
            if not self.auto_y_right.get():
                ymin = _safe_float(self.y_right_min.get())
                ymax = _safe_float(self.y_right_max.get())
                if ymin is not None and ymax is not None:
                    ax.set_ylim(ymin, ymax)

    def _plot_series(self, ax, x, y, label):
        st = self.series_style.get(label, {})
        color = st.get("color", None)
        marker = st.get("marker", MARKER_DEFAULT)
        lw = max(0.1, float(st.get("linewidth", self.linewidth.get())))
        ms = max(0.1, float(st.get("markersize", self.markersize.get())))
        kind = self.plot_kind.get()
        if kind == "scatter":
            return ax.scatter(x, y, s=ms, label=label, color=color, marker=marker)
        elif kind == "bar":
            return ax.bar(x, y, label=label, color=color)
        elif kind == "area":
            return ax.fill_between(x, 0, y, label=label, alpha=0.35, color=color)
        elif kind == "step":
            (h,) = ax.step(x, y, label=label, where="mid", color=color)
            return h
        else:  # line
            (h,) = ax.plot(
                x,
                y,
                linestyle=LINESTYLE_DEFAULT,
                marker=marker,
                linewidth=lw,
                markersize=ms,
                label=label,
                color=color,
            )
            return h

    def _transform_series(self, s: pd.Series, label: str) -> pd.Series:
        st = self.series_style.get(label, {})
        if not st:
            return s
        try:
            scale = float(st.get("scale", 1.0))
            offset = float(st.get("offset", 0.0))
            win = int(st.get("rolling", 0))
        except Exception:
            scale, offset, win = 1.0, 0.0, 0
        s = s * scale + offset
        if win and win > 1:
            s = s.rolling(win, min_periods=1).mean()
        return s

    def _apply_references(self, ax):
        for ref in self.references:
            t = ref.get("type")
            if t == "hline":
                y = _safe_float(ref.get("y"))
                if y is not None:
                    ax.axhline(
                        y,
                        linestyle="--",
                        alpha=float(ref.get("alpha", 0.5)),
                        color=ref.get("color", None),
                        label=ref.get("label", None),
                    )
            elif t == "vline":
                x = _safe_float(ref.get("x"))
                if x is not None:
                    ax.axvline(
                        x,
                        linestyle="--",
                        alpha=float(ref.get("alpha", 0.5)),
                        color=ref.get("color", None),
                        label=ref.get("label", None),
                    )
            elif t == "band_y":
                y0 = _safe_float(ref.get("y0"))
                y1 = _safe_float(ref.get("y1"))
                if y0 is not None and y1 is not None:
                    ax.axhspan(
                        y0,
                        y1,
                        alpha=float(ref.get("alpha", 0.15)),
                        color=ref.get("color", None),
                        label=ref.get("label", None),
                    )
            elif t == "band_x":
                x0 = _safe_float(ref.get("x0"))
                x1 = _safe_float(ref.get("x1"))
                if x0 is not None and x1 is not None:
                    ax.axvspan(
                        x0,
                        x1,
                        alpha=float(ref.get("alpha", 0.15)),
                        color=ref.get("color", None),
                        label=ref.get("label", None),
                    )

    # ---------- Live preview ----------
    def draw_preview(self):
        self._preview_after_id = None
        self.preview_fig.clf()
        _apply_theme(self.theme.get())

        if self.df is None:
            ax = self.preview_fig.add_subplot(111)
            ax.text(
                0.5,
                0.5,
                "Load data on the Data tab.",
                ha="center",
                va="center",
                fontsize=12,
            )
            self.preview_canvas.draw()
            return

        xcol = self.cb_x.get()
        ycols = self._get_selected_y()
        if not xcol or not ycols:
            ax = self.preview_fig.add_subplot(111)
            ax.text(
                0.5,
                0.5,
                "Select X and at least one Y.",
                ha="center",
                va="center",
                fontsize=12,
            )
            self.preview_canvas.draw()
            return

        df = self.df
        q = self.filter_query.get().strip()
        if q:
            try:
                df = df.query(q)
            except Exception:
                pass
        if self.decimate_preview.get() and len(df) > 5000:
            df = df.iloc[:: max(1, len(df) // 5000)].copy()

        x = pd.to_numeric(df[xcol], errors="coerce")
        data = {col: pd.to_numeric(df[col], errors="coerce") for col in ycols}
        ptype = self.plot_type.get()
        legend_loc = self.legend_loc.get()
        legend_outside = legend_loc == "outside bottom"
        right_series = self.right_axis_series.get() if ptype.endswith("right") else None

        if ptype in ("single", "single_right"):
            ax = self.preview_fig.add_subplot(111)
            ax.grid(
                self.show_grid.get(), which="both", axis="both"
            )  # toggle grid on/off
            for col in ycols:
                st = self.series_style.get(col, {})
                if st.get("visible", True) is False:
                    continue
                y_plot = self._transform_series(data[col], col)
                self._plot_series(ax, x, y_plot, col)

            yc = self.yerr_col.get()
            if yc and yc != "None" and yc in df.columns:
                yerr = pd.to_numeric(df[yc], errors="coerce")
                try:
                    ax.errorbar(
                        x,
                        self._transform_series(data[ycols[0]], ycols[0]),
                        yerr=yerr,
                        fmt="none",
                        alpha=0.5,
                        capsize=3,
                    )
                except Exception:
                    pass

            if ptype.endswith("right"):
                if (
                    right_series
                    and right_series != "None"
                    and right_series in df.columns
                ):
                    ax2 = ax.twinx()
                    ax2.grid(False)
                    self._plot_series(
                        ax2,
                        x,
                        pd.to_numeric(df[right_series], errors="coerce"),
                        right_series,
                    )
                    self._apply_limits(ax2, which="right", set_x=False)
                    self._apply_ticks(ax2, set_x=False)
                    _apply_formatter_to_axis(ax2.yaxis, self.y_right_format.get())
                    if self.log_y_right.get():
                        ax2.set_yscale("log")
                    ax2.set_ylabel(
                        _resolve_right_label(self.y_right_label.get(), right_series),
                        fontsize=LABEL_FONTSIZE,
                        rotation=-90,
                        labelpad=12,
                    )
                else:
                    ax.text(
                        0.5,
                        0.1,
                        "Select a Right-Y series.",
                        ha="center",
                        transform=ax.transAxes,
                        fontsize=10,
                    )

            # Labels (use custom if provided)
            ax.set_xlabel(self.x_label.get() or xcol, fontsize=LABEL_FONTSIZE)
            left_label_default = (
                ", ".join([c for c in ycols if c != right_series]) or "Y"
            )
            ax.set_ylabel(
                self.y_left_label.get() or left_label_default, fontsize=LABEL_FONTSIZE
            )

            self._apply_ticks(ax)
            self._apply_limits(ax, which="left")
            _apply_formatter_to_axis(ax.yaxis, self.y_left_format.get())
            self._apply_references(ax)

            if self.log_x.get():
                ax.set_xscale("log")
            if self.log_y_left.get():
                ax.set_yscale("log")
            for lbl in ax.get_xticklabels():
                lbl.set_rotation(self.x_tick_rotation.get())

            ax.set_title(
                self.title_text.get() or "Preview", fontsize=SUBPLOT_TITLE_FONTSIZE
            )
            self.preview_fig.suptitle(
                self.suptitle_text.get(), fontsize=SUPTITLE_FONTSIZE, y=SUPTITLE_Y
            )

        elif ptype in ("grid", "grid_right"):
            n = len(ycols)
            ncols = max(1, int(self.ncols.get()))
            nrows = int(np.ceil(n / ncols))
            axs = self.preview_fig.subplots(nrows=nrows, ncols=ncols)
            axs = np.array(axs).reshape(-1)
            right_series = (
                self.right_axis_series.get() if ptype.endswith("right") else None
            )

            for i, col in enumerate(ycols):
                ax = axs[i]
                ax.grid(self.show_grid.get())
                st = self.series_style.get(col, {})
                if st.get("visible", True) is False:
                    ax.set_visible(True)
                else:
                    y_plot = self._transform_series(data[col], col)
                    self._plot_series(ax, x, y_plot, col)

                yc = self.yerr_col.get()
                if yc and yc != "None" and yc in df.columns:
                    yerr = pd.to_numeric(df[yc], errors="coerce")
                    try:
                        ax.errorbar(
                            x,
                            self._transform_series(data[col], col),
                            yerr=yerr,
                            fmt="none",
                            alpha=0.4,
                            capsize=2,
                        )
                    except Exception:
                        pass

                if right_series and right_series in df.columns:
                    ax2 = ax.twinx()
                    ax2.grid(False)
                    self._plot_series(
                        ax2,
                        x,
                        pd.to_numeric(df[right_series], errors="coerce"),
                        right_series,
                    )
                    self._apply_limits(ax2, which="right", set_x=False)
                    self._apply_ticks(ax2, set_x=False)
                    _apply_formatter_to_axis(ax2.yaxis, self.y_right_format.get())
                    if self.log_y_right.get():
                        ax2.set_yscale("log")
                    # Right axis label
                    ax2.set_ylabel(
                        _resolve_right_label(self.y_right_label.get(), right_series),
                        fontsize=LABEL_FONTSIZE,
                        rotation=-90,
                        labelpad=12,
                    )

                ax.set_title(col, fontsize=SUBPLOT_TITLE_FONTSIZE)
                ax.set_xlabel(self.x_label.get() or xcol, fontsize=LABEL_FONTSIZE)
                ax.set_ylabel(self.y_left_label.get() or col, fontsize=LABEL_FONTSIZE)
                self._apply_ticks(ax)
                self._apply_limits(ax, which="left")
                _apply_formatter_to_axis(ax.yaxis, self.y_left_format.get())
                self._apply_references(ax)
                if legend_loc != "outside bottom":
                    ax.legend(loc="best", fontsize=LABEL_FONTSIZE - 1)

            for j in range(n, len(axs)):
                axs[j].set_visible(False)
            if legend_loc == "outside bottom":
                # Gather unique handles/labels from all subplots (and any right twins)
                handles, labels = [], []
                twins = []
                for i, col in enumerate(ycols):
                    # collect from left axes
                    h, l = axs[i].get_legend_handles_labels()
                    handles += h
                    labels += l
                    # collect from right axes if created
                    for child in axs[i].get_children():
                        # no reliable handle, so instead track twins explicitly during creation
                        pass
                # If you created ax2 twins in the loop, append them to `twins` there; then:
                for t in twins:
                    try:
                        h2, l2 = t.get_legend_handles_labels()
                        handles += h2
                        labels += l2
                    except Exception:
                        pass

                seen, H, L = set(), [], []
                for h, l in zip(handles, labels):
                    if l and l not in seen:
                        seen.add(l)
                        H.append(h)
                        L.append(l)

                if H:
                    self.preview_fig.legend(
                        H,
                        L,
                        loc="upper center",
                        bbox_to_anchor=(0.5, 0.02),  # inside reserved bottom margin
                        ncol=min(3, len(L)),
                        frameon=True,
                        fancybox=True,
                        shadow=True,
                        borderaxespad=0.0,
                    )
                # Make room under subplots
                self.preview_fig.subplots_adjust(bottom=0.22)
            if legend_loc == "outside bottom":
                # Gather unique handles/labels from all subplots (and any right twins)
                handles, labels = [], []
                twins = []
                for i, col in enumerate(ycols):
                    # collect from left axes
                    h, l = axs[i].get_legend_handles_labels()
                    handles += h
                    labels += l
                    # collect from right axes if created
                    for child in axs[i].get_children():
                        # no reliable handle, so instead track twins explicitly during creation
                        pass
                # If you created ax2 twins in the loop, append them to `twins` there; then:
                for t in twins:
                    try:
                        h2, l2 = t.get_legend_handles_labels()
                        handles += h2
                        labels += l2
                    except Exception:
                        pass

                seen, H, L = set(), [], []
                for h, l in zip(handles, labels):
                    if l and l not in seen:
                        seen.add(l)
                        H.append(h)
                        L.append(l)

                if H:
                    self.preview_fig.legend(
                        H,
                        L,
                        loc="upper center",
                        bbox_to_anchor=(0.5, 0.02),  # inside reserved bottom margin
                        ncol=min(3, len(L)),
                        frameon=True,
                        fancybox=True,
                        shadow=True,
                        borderaxespad=0.0,
                    )
                # Make room under subplots
                self.preview_fig.subplots_adjust(bottom=0.22)

            self.preview_fig.tight_layout(rect=[0, 0, 1, 0.94])
            self.preview_fig.suptitle(
                self.suptitle_text.get() or self.title_text.get(),
                fontsize=SUPTITLE_FONTSIZE,
                y=SUPTITLE_Y,
            )
            if legend_loc == "outside bottom":
                self.preview_fig.subplots_adjust(bottom=0.22)
                handles, labels = self._collect_legend_entries(self.preview_fig)
                if handles:
                    self.preview_fig.legend(
                        handles,
                        labels,
                        loc="upper center",
                        bbox_to_anchor=(0.5, 0.02),
                        ncol=min(3, len(labels)),
                        frameon=True,
                        fancybox=True,
                        shadow=True,
                        borderaxespad=0.0,
                    )

            try:
                txt = []
                for col in ycols[:6]:
                    s = pd.to_numeric(df[col], errors="coerce")
                    txt.append(
                        f"{col}: min={s.min():.3g}, mean={s.mean():.3g}, max={s.max():.3g}"
                    )
                self.stats_text.delete("1.0", "end")
                self.stats_text.insert("1.0", "\n".join(txt) or "—")
            except Exception:
                pass

        else:  # facet / facet_right
            facet_col = self.cb_facet.get()
            if facet_col in (None, "", "None"):
                ax = self.preview_fig.add_subplot(111)
                ax.text(
                    0.5,
                    0.5,
                    "Choose a Facet column (category) for facet plots.",
                    ha="center",
                    va="center",
                    fontsize=12,
                )
                self.preview_canvas.draw()
                return
            cats = pd.Series(df[facet_col]).astype("category")
            levels = list(cats.cat.categories)
            if not levels:
                ax = self.preview_fig.add_subplot(111)
                ax.text(
                    0.5,
                    0.5,
                    f"Column '{facet_col}' has no categories.",
                    ha="center",
                    va="center",
                    fontsize=12,
                )
                self.preview_canvas.draw()
                return

            n = len(levels)
            ncols = max(1, int(self.ncols.get()))
            nrows = int(np.ceil(n / ncols))
            axs = self.preview_fig.subplots(nrows=nrows, ncols=ncols)
            axs = np.array(axs).reshape(-1)
            right_series = (
                self.right_axis_series.get() if ptype.endswith("right") else None
            )

            for i, lvl in enumerate(levels):
                ax = axs[i]
                ax.grid(self.show_grid.get())
                mask = cats == lvl
                x_sub = pd.to_numeric(df.loc[mask, xcol], errors="coerce")
                for col in ycols:
                    y_sub = pd.to_numeric(df.loc[mask, col], errors="coerce")
                    st = self.series_style.get(col, {})
                    if st.get("visible", True) is not False:
                        self._plot_series(
                            ax, x_sub, self._transform_series(y_sub, col), col
                        )

                if right_series and right_series in df.columns:
                    ax2 = ax.twinx()
                    ax2.grid(False)
                    y_r = pd.to_numeric(df.loc[mask, right_series], errors="coerce")
                    self._plot_series(ax2, x_sub, y_r, right_series)
                    self._apply_limits(ax2, which="right", set_x=False)
                    self._apply_ticks(ax2, set_x=False)
                    _apply_formatter_to_axis(ax2.yaxis, self.y_right_format.get())
                    if self.log_y_right.get():
                        ax2.set_yscale("log")
                    ax2.set_ylabel(
                        _resolve_right_label(self.y_right_label.get(), right_series),
                        fontsize=LABEL_FONTSIZE,
                        rotation=-90,
                        labelpad=12,
                    )

                ax.set_title(f"{facet_col} = {lvl}", fontsize=SUBPLOT_TITLE_FONTSIZE)
                ax.set_xlabel(self.x_label.get() or xcol, fontsize=LABEL_FONTSIZE)
                ax.set_ylabel(
                    self.y_left_label.get() or (", ".join(ycols)),
                    fontsize=LABEL_FONTSIZE,
                )
                self._apply_ticks(ax)
                self._apply_limits(ax, which="left")
                _apply_formatter_to_axis(ax.yaxis, self.y_left_format.get())
                self._apply_references(ax)

                if self.log_x.get():
                    ax.set_xscale("log")
                if self.log_y_left.get():
                    ax.set_yscale("log")
                for lbl in ax.get_xticklabels():
                    lbl.set_rotation(self.x_tick_rotation.get())
                if legend_loc != "outside bottom":
                    ax.legend(loc="best", fontsize=LABEL_FONTSIZE - 1)

            for j in range(n, len(axs)):
                axs[j].set_visible(False)
            self.preview_fig.tight_layout(rect=[0, 0, 1, 0.94])
            self.preview_fig.suptitle(
                self.suptitle_text.get() or self.title_text.get(),
                fontsize=SUPTITLE_FONTSIZE,
                y=SUPTITLE_Y,
            )
            if legend_outside:
                # Make room for the outside legend and place a single figure-level legend
                self.preview_fig.subplots_adjust(bottom=0.20)
                handles, labels = self._collect_legend_entries(self.preview_fig)
                if handles:
                    self.preview_fig.legend(
                        handles,
                        labels,
                        loc="lower center",
                        ncol=min(6, len(labels)),
                        fontsize=LABEL_FONTSIZE - 1,
                        frameon=False,
                    )

        if mplcursors is not None:
            try:
                mplcursors.cursor(self.preview_fig, hover=True)
            except Exception:
                pass
        self.preview_canvas.draw()

    # ---------- Limits/ticks state ----------
    def _sync_range_states(self):
        for widget, flag in [(self.e_xmin, self.auto_x), (self.e_xmax, self.auto_x)]:
            widget.configure(state="disabled" if flag.get() else "normal")
        for widget, flag in [
            (self.e_ylmin, self.auto_y_left),
            (self.e_ylmax, self.auto_y_left),
        ]:
            widget.configure(state="disabled" if flag.get() else "normal")
        for widget, flag in [
            (self.e_yrmin, self.auto_y_right),
            (self.e_yrmax, self.auto_y_right),
        ]:
            widget.configure(state="disabled" if flag.get() else "normal")
        self._schedule_preview()

    # ---------- Full render/export ----------
    def render_plots(self):
        if self.df is None:
            messagebox.showerror("No data", "Load a file/sheet first.")
            return
        xcol = self.cb_x.get()
        ycols = self._get_selected_y()
        if not xcol:
            messagebox.showerror("Missing X", "Select an X column.")
            return
        if not ycols:
            messagebox.showerror("Missing Y", "Select at least one Y column.")
            return

        _apply_theme(self.theme.get())
        df_full = self.df
        q = self.filter_query.get().strip()
        if q:
            try:
                df_full = df_full.query(q)
            except Exception:
                pass

        ptype = self.plot_type.get()
        x = pd.to_numeric(df_full[xcol], errors="coerce")
        data = {col: pd.to_numeric(df_full[col], errors="coerce") for col in ycols}
        legend_loc = self.legend_loc.get()
        legend_outside = legend_loc == "outside bottom"

        figs = []

        def apply_common(ax):
            # Labels (use custom if provided)
            ax.set_xlabel(self.x_label.get() or xcol, fontsize=LABEL_FONTSIZE)
            self._apply_ticks(ax)
            self._apply_limits(ax, which="left")
            if self.log_x.get():
                ax.set_xscale("log")
            if self.log_y_left.get():
                ax.set_yscale("log")
            for lbl in ax.get_xticklabels():
                lbl.set_rotation(self.x_tick_rotation.get())
            if legend_loc != "outside bottom":
                ax.legend(
                    loc=legend_loc if ptype.startswith("single") else "best",
                    fontsize=LABEL_FONTSIZE,
                    ncol=min(3, len(ycols)),
                )

        if ptype in ("single", "single_right"):
            fig, ax = plt.subplots(figsize=(11, 8.5))
            fig.subplots_adjust(left=0.075, right=0.92, bottom=0.14, top=0.91)
            ax.grid(self.show_grid.get())
            right_choice = None
            for col in ycols:
                st = self.series_style.get(col, {})
                if st.get("visible", True) is False:
                    continue
                y_plot = self._transform_series(data[col], col)
                self._plot_series(ax, x, y_plot, col)

            if ptype.endswith("right"):
                right_choice = self.right_axis_series.get()
                if (
                    right_choice
                    and right_choice != "None"
                    and right_choice in self.df.columns
                ):
                    ax2 = ax.twinx()
                    self._plot_series(
                        ax2,
                        x,
                        pd.to_numeric(self.df[right_choice], errors="coerce"),
                        right_choice,
                    )
                    self._apply_limits(ax2, which="right", set_x=False)
                    self._apply_ticks(ax2, set_x=False)
                    _apply_formatter_to_axis(ax2.yaxis, self.y_right_format.get())
                    if self.log_y_right.get():
                        ax2.set_yscale("log")
                    ax2.set_ylabel(
                        _resolve_right_label(self.y_right_label.get(), right_choice),
                        fontsize=LABEL_FONTSIZE,
                        rotation=-90,
                        labelpad=15,
                    )

                else:
                    messagebox.showwarning(
                        "Right Y", "Select a valid series for the right Y-axis."
                    )

            left_label_default = (
                ", ".join([c for c in ycols if c != right_choice]) or "Y"
            )
            ax.set_ylabel(
                self.y_left_label.get() or left_label_default, fontsize=LABEL_FONTSIZE
            )
            ax.set_title(self.title_text.get(), fontsize=SUBPLOT_TITLE_FONTSIZE)
            fig.suptitle(
                self.suptitle_text.get(), fontsize=SUPTITLE_FONTSIZE, y=SUPTITLE_Y
            )
            self._apply_references(ax)
            apply_common(ax)
            _apply_formatter_to_axis(ax.yaxis, self.y_left_format.get())

            if legend_loc == "outside bottom":
                handles, labels = ax.get_legend_handles_labels()
                try:
                    if "ax2" in locals() and ax2 is not None:
                        h2, l2 = ax2.get_legend_handles_labels()
                        handles += h2
                        labels += l2
                except Exception:
                    pass

                seen, H, L = set(), [], []
                for h, l in zip(handles, labels):
                    if l and l not in seen:
                        seen.add(l)
                        H.append(h)
                        L.append(l)

                if H:
                    ax.legend(
                        H,
                        L,
                        loc="upper center",
                        bbox_to_anchor=(0.5, -0.26),
                        ncol=min(3, len(L)),
                        frameon=True,
                        fancybox=True,
                        shadow=True,
                        borderaxespad=0.0,
                    )
                fig.subplots_adjust(bottom=0.28)

            _make_legends_draggable(fig)

            figs.append(fig)

        elif ptype in ("grid", "grid_right"):
            n = len(ycols)
            ncols = max(1, int(self.ncols.get()))
            nrows = int(np.ceil(n / ncols))
            fig, axes = plt.subplots(
                nrows=nrows, ncols=ncols, figsize=(12, max(6, 3 * nrows))
            )
            axes = np.array(axes).reshape(-1)
            fig.subplots_adjust(
                left=0.07, right=0.95, bottom=0.08, top=0.9, hspace=0.35, wspace=0.25
            )
            right_choice = (
                self.right_axis_series.get() if ptype.endswith("right") else None
            )
            right_twins = []
            for i, col in enumerate(ycols):
                ax = axes[i]
                ax.grid(self.show_grid.get())
                st = self.series_style.get(col, {})
                if st.get("visible", True) is not False:
                    y_plot = self._transform_series(data[col], col)
                    self._plot_series(ax, x, y_plot, col)
                if right_choice and right_choice in self.df.columns:
                    ax2 = ax.twinx()
                    right_twins.append(ax2)
                    self._plot_series(
                        ax2,
                        x,
                        pd.to_numeric(self.df[right_choice], errors="coerce"),
                        right_choice,
                    )
                    self._apply_limits(ax2, which="right", set_x=False)
                    self._apply_ticks(ax2, set_x=False)
                    _apply_formatter_to_axis(ax2.yaxis, self.y_right_format.get())
                    if self.log_y_right.get():
                        ax2.set_yscale("log")
                    ax2.set_ylabel(
                        _resolve_right_label(self.y_right_label.get(), right_choice),
                        fontsize=LABEL_FONTSIZE,
                        rotation=-90,
                        labelpad=12,
                    )
            ax.set_title(col, fontsize=SUBPLOT_TITLE_FONTSIZE)
            ax.set_ylabel(self.y_left_label.get() or col, fontsize=LABEL_FONTSIZE)
            self._apply_references(ax)
            apply_common(ax)
            _apply_formatter_to_axis(ax.yaxis, self.y_left_format.get())

            for j in range(len(ycols), len(axes)):
                axes[j].set_visible(False)
            if legend_loc == "outside bottom":
                handles, labels = [], []
                # collect from each subplot
                for ax_i in axes:
                    h, l = ax_i.get_legend_handles_labels()
                    handles += h
                    labels += l
                # collect from right-side twins if any
                for t in right_twins:
                    try:
                        h2, l2 = t.get_legend_handles_labels()
                        handles += h2
                        labels += l2
                    except Exception:
                        pass

                # de-duplicate while preserving order
                seen, H, L = set(), [], []
                for h, l in zip(handles, labels):
                    if l and l not in seen:
                        seen.add(l)
                        H.append(h)
                        L.append(l)

                if H:
                    fig.legend(
                        H,
                        L,
                        loc="upper center",
                        bbox_to_anchor=(0.5, 0.02),
                        ncol=min(3, len(L)),
                        frameon=True,
                        fancybox=True,
                        shadow=True,
                        borderaxespad=0.0,
                    )
                fig.subplots_adjust(bottom=0.22)

            _make_legends_draggable(fig)

            fig.suptitle(
                self.suptitle_text.get() or self.title_text.get(),
                fontsize=SUPTITLE_FONTSIZE,
                y=SUPTITLE_Y,
            )
            figs.append(fig)

        else:  # facet / facet_right
            facet_col = self.cb_facet.get()
            if facet_col in (None, "", "None"):
                messagebox.showerror(
                    "Facet missing", "Choose a Facet column for facet plots."
                )
                return
            cats = pd.Series(self.df[facet_col]).astype("category")
            levels = list(cats.cat.categories)
            k = int(self.top_k_facets.get() or 0)
            if k > 0 and len(levels) > k:
                counts = cats.value_counts().to_dict()
                levels.sort(key=lambda v: counts.get(v, 0), reverse=True)
                levels = levels[:k]
            if not levels:
                messagebox.showerror(
                    "No categories", f"Column '{facet_col}' has no categories."
                )
                return

            n = len(levels)
            ncols = max(1, int(self.ncols.get()))
            nrows = int(np.ceil(n / ncols))
            fig, axes = plt.subplots(
                nrows=nrows, ncols=ncols, figsize=(13, max(6, 3 * nrows))
            )
            axes = np.array(axes).reshape(-1)
            fig.subplots_adjust(
                left=0.07, right=0.95, bottom=0.08, top=0.9, hspace=0.35, wspace=0.25
            )
            right_choice = (
                self.right_axis_series.get() if ptype.endswith("right") else None
            )
            right_twins = []

            for i, lvl in enumerate(levels):
                ax = axes[i]
                ax.grid(self.show_grid.get())
                mask = cats == lvl
                x_sub = pd.to_numeric(self.df.loc[mask, xcol], errors="coerce")
                for col in ycols:
                    y_sub = pd.to_numeric(self.df.loc[mask, col], errors="coerce")
                    st = self.series_style.get(col, {})
                    if st.get("visible", True) is not False:
                        self._plot_series(
                            ax, x_sub, self._transform_series(y_sub, col), col
                        )
                if right_choice and right_choice in self.df.columns:
                    ax2 = ax.twinx()
                    right_twins.append(ax2)
                    y_r = pd.to_numeric(
                        self.df.loc[mask, right_choice], errors="coerce"
                    )
                    self._plot_series(ax2, x_sub, y_r, right_choice)
                    self._apply_limits(ax2, which="right", set_x=False)
                    self._apply_ticks(ax2, set_x=False)
                    _apply_formatter_to_axis(ax2.yaxis, self.y_right_format.get())
                    if self.log_y_right.get():
                        ax2.set_yscale("log")
                    ax2.set_ylabel(
                        _resolve_right_label(self.y_right_label.get(), right_choice),
                        fontsize=LABEL_FONTSIZE,
                        rotation=-90,
                        labelpad=12,
                    )
                ax.set_title(f"{facet_col} = {lvl}", fontsize=SUBPLOT_TITLE_FONTSIZE)
                ax.set.ylabel(
                    self.y_left_label.get() or (", ".join(ycols)),
                    fontsize=LABEL_FONTSIZE,
                )

                apply_common(ax)
                _apply_formatter_to_axis(ax.yaxis, self.y_left_format.get())
                self._apply_references(ax)
            for j in range(n, len(axes)):
                axes[j].set_visible(False)
            if legend_loc == "outside bottom":
                handles, labels = [], []
                for ax_i in axes:
                    h, l = ax_i.get_legend_handles_labels()
                    handles += h
                    labels += l
                for t in right_twins:
                    try:
                        h2, l2 = t.get_legend_handles_labels()
                        handles += h2
                        labels += l2
                    except Exception:
                        pass

                seen, H, L = set(), [], []
                for h, l in zip(handles, labels):
                    if l and l not in seen:
                        seen.add(l)
                        H.append(h)
                        L.append(l)

                if H:
                    fig.legend(
                        H,
                        L,
                        loc="upper center",
                        bbox_to_anchor=(0.5, 0.02),
                        ncol=min(3, len(L)),
                        frameon=True,
                        fancybox=True,
                        shadow=True,
                        borderaxespad=0.0,
                    )
                fig.subplots_adjust(bottom=0.22)

            _make_legends_draggable(fig)

            fig.suptitle(
                self.suptitle_text.get() or self.title_text.get(),
                fontsize=SUPTITLE_FONTSIZE,
                y=SUPTITLE_Y,
            )
            figs.append(fig)

        if figs:
            self._display_rendered_figures(figs)
        self._save_settings()

    def export_plots(self):
        if not plt.get_fignums():
            messagebox.showerror("Nothing to export", "Render a plot first.")
            return
        outdir = filedialog.askdirectory(title="Choose export folder")
        if not outdir:
            return
        base = self.title_text.get().strip() or "plot"
        base = (
            "".join(ch for ch in base if ch.isalnum() or ch in (" ", "_", "-"))
            .strip()
            .replace(" ", "_")
        )
        for i, num in enumerate(plt.get_fignums(), 1):
            fig = plt.figure(num)
            png = os.path.join(outdir, f"{base}_{i:02d}.png")
            svg = os.path.join(outdir, f"{base}_{i:02d}.svg")
            pdf = os.path.join(outdir, f"{base}_{i:02d}.pdf")
            try:
                fig.savefig(png, dpi=EXPORT_DPI, bbox_inches="tight")
                fig.savefig(svg, bbox_inches="tight")
                fig.savefig(pdf, dpi=EXPORT_DPI, bbox_inches="tight")
            except Exception as e:
                messagebox.showerror("Export error", f"Could not save figure {i}: {e}")
                return
        messagebox.showinfo(
            "Export complete", f"Saved {len(plt.get_fignums())} figure(s) to:\n{outdir}"
        )


if __name__ == "__main__":
    app = GeneralPlotter()
    try:
        app.draw_preview()
    except Exception:
        pass
    app.mainloop()
    try:
        plt.close("all")
    except Exception:
        pass
    sys.exit(0)
